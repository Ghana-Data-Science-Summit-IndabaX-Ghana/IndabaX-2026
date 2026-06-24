"""
Build the 'Automating Data Pipelines with Python' workshop deck.
Front End Masters style: code-first, clear mental models, premium dark design.

Palette: agrivoltaics -> solar gold + crop green on deep forest-charcoal.
Code snippets mirror the REAL pipeline in /Users/Apple/Practice/agrivoltic-pipeline.
Output: Workshop_Slide_Deck.pptx
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Design tokens
# ----------------------------------------------------------------------------
BG       = RGBColor(0x0E, 0x1A, 0x13)
PANEL    = RGBColor(0x17, 0x27, 0x1D)
PANEL2   = RGBColor(0x1E, 0x33, 0x27)
CODE_BG  = RGBColor(0x0A, 0x14, 0x10)
GREEN    = RGBColor(0x54, 0xD1, 0x7A)
GREEN_D  = RGBColor(0x2F, 0x7A, 0x4B)
GOLD     = RGBColor(0xF7, 0xC9, 0x48)
CORAL    = RGBColor(0xE0, 0x6C, 0x5E)
TEXT     = RGBColor(0xF3, 0xF1, 0xE9)
MUTED    = RGBColor(0x9F, 0xB3, 0xA6)
LINE     = RGBColor(0x2B, 0x40, 0x33)

C_DEF  = RGBColor(0xD7, 0xE0, 0xD5)
C_KW   = RGBColor(0xF7, 0xC9, 0x48)
C_STR  = RGBColor(0x7F, 0xD8, 0x9A)
C_COM  = RGBColor(0x6E, 0x83, 0x77)
C_NUM  = RGBColor(0xF0, 0xA8, 0x68)

HEAD = "Avenir Next"
BODY = "Avenir Next"
MONO = "Menlo"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
TOTAL = 55

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _noline(shape):
    shape.line.fill.background()

def _line(shape, color, w=1.0):
    shape.line.color.rgb = color; shape.line.width = Pt(w)


def new_slide(color=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    _fill(r, color); _noline(r); r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def rrect(slide, x, y, w, h, fill=PANEL, radius=0.08, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        _fill(s, fill)
    if line is None:
        _noline(s)
    else:
        _line(s, line, line_w)
    s.shadow.inherit = False
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s


def oval(slide, x, y, d, fill, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                               Inches(d), Inches(d))
    _fill(s, fill)
    if line is None:
        _noline(s)
    else:
        _line(s, line, line_w)
    s.shadow.inherit = False
    return s


def _spacing(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def text(slide, x, y, w, h, runs, size=14, color=TEXT, bold=False, italic=False,
         font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0,
         space_after=0, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(space_after)
        try:
            p.line_spacing = spacing
        except Exception:
            pass
        if isinstance(para, str):
            para = [(para, {})]
        for txt, opts in para:
            r = p.add_run(); r.text = txt; f = r.font
            f.size = Pt(opts.get("size", size)); f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", italic); f.name = opts.get("font", font)
            f.color.rgb = opts.get("color", color)
            if opts.get("spacing") is not None:
                _spacing(r, opts["spacing"])
    return tb


def pill(slide, x, y, label, fill=GOLD, fg=BG, size=11, bold=True, letter=0.6):
    w = 0.30 + 0.098 * (size / 11) * len(label) + 0.02 * letter * len(label)
    h = 0.30 * size / 11 + 0.06
    s = rrect(slide, x, y, w, h, fill=fill, radius=0.5)
    tf = s.text_frame; tf.word_wrap = False
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = HEAD; r.font.color.rgb = fg
    _spacing(r, letter)
    return w, h


# ---- code panel ------------------------------------------------------------
KEYWORDS = {"def", "return", "if", "elif", "else", "for", "in", "from",
            "import", "raise", "not", "and", "or", "with", "as", "None",
            "True", "False", "lambda", "class", "while", "is", "try",
            "except", "pass", "set", "sorted", "round", "float", "int", "str"}
TOKEN_RE = re.compile(r'(\".*?\"|\'.*?\'|#.*$|\b\w+\b|\s+|[^\w\s])')


def _tok(line):
    if line.lstrip().startswith("#"):
        return [(line, C_COM, True)]
    out = []
    for t in TOKEN_RE.findall(line):
        if not t:
            continue
        if t.startswith("#"):
            out.append((t, C_COM, True))
        elif t[:1] in "\"'":
            out.append((t, C_STR, False))
        elif t in KEYWORDS:
            out.append((t, C_KW, False))
        elif re.fullmatch(r"\d+\.?\d*", t):
            out.append((t, C_NUM, False))
        else:
            out.append((t, C_DEF, False))
    return out


def code_panel(slide, x, y, w, h, filename, code, size=13):
    rrect(slide, x, y, w, h, fill=CODE_BG, radius=0.05, line=LINE, line_w=1.0)
    for i, c in enumerate([CORAL, RGBColor(0xE6, 0xB4, 0x5C), GREEN]):
        oval(slide, x + 0.28 + i * 0.26, y + 0.26, 0.13, c)
    text(slide, x + 1.25, y + 0.16, w - 1.4, 0.32, filename, size=11,
         color=MUTED, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    tb = slide.shapes.add_textbox(Inches(x + 0.34), Inches(y + 0.62),
                                  Inches(w - 0.6), Inches(h - 0.8))
    tf = tb.text_frame; tf.word_wrap = False
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0); p.space_after = Pt(0); p.line_spacing = 1.16
        toks = _tok(ln) if ln.strip() else [(" ", C_DEF, False)]
        for t, col, ital in toks:
            r = p.add_run(); r.text = t; r.font.name = MONO
            r.font.size = Pt(size); r.font.color.rgb = col; r.font.italic = ital


# ----------------------------------------------------------------------------
# Composite components
# ----------------------------------------------------------------------------

def kicker(slide, label, fill=GOLD, x=0.7, y=0.62):
    pill(slide, x, y, label, fill=fill, fg=BG, size=11)


def title_block(slide, title, sub=None, x=0.7, y=1.12, w=11.9, size=32):
    text(slide, x, y, w, 1.0, title, size=size, bold=True, color=TEXT,
         font=HEAD, spacing=1.02)
    if sub:
        text(slide, x + 0.02, y + 0.66, w, 0.5, sub, size=15, color=MUTED, font=BODY)


def footer(slide, label, num):
    text(slide, 0.7, 7.05, 9, 0.3, label, size=9.5, color=MUTED, font=BODY)
    text(slide, 11.0, 7.05, 1.63, 0.3, f"{num} / {TOTAL}", size=9.5,
         color=MUTED, font=BODY, align=PP_ALIGN.RIGHT)


def arrow(slide, x, y, color=GOLD, size=18):
    text(slide, x, y, 0.4, 0.4, "→", size=size, color=color,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)


def step_flow(slide, x, y, w, steps, chip_h=0.62, color=GOLD, fill=PANEL,
              label_size=11):
    n = len(steps); aw = 0.34
    cw = (w - (n - 1) * aw) / n
    cx = x
    for i, (num, label) in enumerate(steps):
        rrect(slide, cx, y, cw, chip_h, fill=fill, radius=0.16, line=LINE, line_w=1.0)
        text(slide, cx + 0.18, y + 0.09, cw - 0.3, 0.25, num, size=11, bold=True,
             color=color, font=HEAD)
        text(slide, cx + 0.18, y + 0.30, cw - 0.3, 0.28, label, size=label_size,
             bold=True, color=TEXT, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            arrow(slide, cx + cw + 0.01, y + chip_h / 2 - 0.2, color=color)
        cx += cw + aw


def card(slide, x, y, w, h, head, lines, accent=GREEN, head_size=15,
         body_size=13):
    rrect(slide, x, y, w, h, fill=PANEL, radius=0.07)
    oval(slide, x + 0.3, y + 0.34, 0.16, accent)
    text(slide, x + 0.6, y + 0.26, w - 0.8, 0.4, head, size=head_size,
         bold=True, color=TEXT, font=HEAD)
    paras = [[(l, {})] for l in lines]
    text(slide, x + 0.32, y + 0.84, w - 0.62, h - 1.0, paras, size=body_size,
         color=MUTED, font=BODY, spacing=1.12, space_after=5)


def concept_cards(slide, items, y=2.55, h=3.5, x=0.7, w=11.9, gap=0.35,
                  head_size=16, body_size=13):
    n = len(items); cw = (w - (n - 1) * gap) / n
    cx = x
    for head, lines, accent in items:
        card(slide, cx, y, cw, h, head, lines, accent=accent,
             head_size=head_size, body_size=body_size)
        cx += cw + gap


def bottom_note(slide, runs, y=6.35, accent=GOLD):
    rrect(slide, 0.7, y, 11.9, 0.62, fill=PANEL, radius=0.1)
    if isinstance(runs, str):
        runs = [(runs, {})]
    text(slide, 1.0, y, 11.3, 0.62, [runs], size=14, color=TEXT, font=BODY,
         anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)


def numbered_circle(slide, x, y, n, d=0.5, color=GREEN):
    c = oval(slide, x, y, d, PANEL2, line=color, line_w=1.5)
    tf = c.text_frame; tf.word_wrap = False
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = color; r.font.name = HEAD


def check_mark(slide, x, y):
    c = oval(slide, x, y, 0.34, GREEN_D)
    tf = c.text_frame; tf.word_wrap = False
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓"; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = TEXT; r.font.name = HEAD


def data_table(slide, x, y, w, headers, rows, col_fr, row_h=0.5, fs=13):
    total = sum(col_fr)
    widths = [w * f / total for f in col_fr]
    # header
    rrect(slide, x, y, w, row_h, fill=PANEL2, radius=0.04)
    cx = x
    for head, cw in zip(headers, widths):
        text(slide, cx + 0.18, y, cw - 0.2, row_h, head, size=fs - 1, bold=True,
             color=GREEN, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw
    # rows
    ry = y + row_h + 0.06
    for r_i, row in enumerate(rows):
        if r_i % 2 == 0:
            rrect(slide, x, ry, w, row_h, fill=PANEL, radius=0.04)
        cx = x
        for c_i, (val, cw) in enumerate(zip(row, widths)):
            col = TEXT if c_i == 0 else MUTED
            text(slide, cx + 0.18, ry, cw - 0.2, row_h, str(val), size=fs,
                 color=col, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
            cx += cw
        ry += row_h + 0.06


# ----------------------------------------------------------------------------
# Reusable full slides
# ----------------------------------------------------------------------------

def divider_slide(num, klabel, title, sub, leave_with, flabel, page):
    s = new_slide(PANEL)
    text(s, 0.68, 1.25, 6, 2.2, num, size=150, bold=True, color=GREEN, font=HEAD)
    kicker(s, klabel, fill=GOLD, x=0.78, y=3.72)
    text(s, 0.7, 4.18, 11.4, 1.0, title, size=46, bold=True, color=TEXT, font=HEAD)
    text(s, 0.72, 5.28, 9.8, 0.8, sub, size=17, color=MUTED, font=BODY, spacing=1.2)
    text(s, 0.72, 6.2, 11.4, 0.5,
         [[("You leave with   ", {"color": MUTED}),
           (leave_with, {"color": TEXT, "bold": True})]], size=14, font=BODY)
    footer(s, flabel, page)
    return s


def exercise_slide(title, sub, minutes, steps, workin, flabel, page):
    s = new_slide()
    kicker(s, "YOUR TURN", fill=GREEN)
    title_block(s, title, sub)
    rrect(s, 10.3, 1.12, 2.3, 1.35, fill=PANEL, radius=0.1, line=GOLD, line_w=1.5)
    text(s, 10.3, 1.32, 2.3, 0.7, str(minutes), size=46, bold=True, color=GOLD,
         font=HEAD, align=PP_ALIGN.CENTER)
    text(s, 10.3, 2.04, 2.3, 0.3, "MINUTES", size=11, bold=True, color=MUTED,
         font=HEAD, align=PP_ALIGN.CENTER)
    n = len(steps)
    y0, step_gap = 2.62, min(0.86, (5.9 - 2.62) / max(n, 1))
    for i, runs in enumerate(steps):
        yy = y0 + i * step_gap
        numbered_circle(s, 0.7, yy, i + 1)
        text(s, 1.45, yy + 0.02, 8.4, 0.5, [runs], size=15.5, color=TEXT,
             font=BODY, anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, 0.7, 6.15, 11.9, 0.62, fill=PANEL, radius=0.12)
    text(s, 1.0, 6.15, 11.3, 0.62,
         [[("WORK IN   ", {"color": MUTED, "bold": True, "size": 12}),
           (workin, {"font": MONO, "color": GOLD, "size": 13.5})]],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, flabel, page)
    return s


def checkpoint_slide(title, sub, checks, flabel, page):
    s = new_slide(PANEL)
    kicker(s, "CHECKPOINT", fill=GREEN, x=0.7, y=0.7)
    text(s, 0.7, 1.25, 11.9, 1.0, title, size=38, bold=True, color=TEXT, font=HEAD)
    text(s, 0.72, 2.12, 11, 0.5, sub, size=16, color=MUTED, font=BODY)
    n = len(checks)
    y0 = 3.05
    gap = min(0.82, (6.6 - y0) / max(n, 1))
    for i, c in enumerate(checks):
        yy = y0 + i * gap
        check_mark(s, 0.9, yy)
        text(s, 1.42, yy - 0.02, 11.0, 0.5, c, size=14.5, color=TEXT, font=BODY,
             anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    footer(s, flabel, page)
    return s


def code_with_notes(page, kick, kfill, title, sub, filename, code, notes_head,
                    notes, flabel, code_w=7.35, csize=13.5):
    s = new_slide()
    kicker(s, kick, fill=kfill)
    title_block(s, title, sub)
    code_panel(s, 0.7, 2.42, code_w, 4.1, filename, code, size=csize)
    nx = 0.7 + code_w + 0.55
    text(s, nx, 2.52, 12.6 - nx, 0.4,
         [[(notes_head, {"color": GREEN, "bold": True, "size": 12, "spacing": 1.2})]],
         font=HEAD)
    yy = 3.02
    for h, b in notes:
        text(s, nx, yy, 12.62 - nx, 0.4, h, size=15, bold=True, color=TEXT, font=HEAD)
        text(s, nx, yy + 0.33, 12.62 - nx, 0.8, b, size=12.5, color=MUTED,
             font=BODY, spacing=1.15)
        yy += 1.18
    footer(s, flabel, page)
    return s


F = "AUTOMATING DATA PIPELINES"

# ============================================================================
# OPENING
# ============================================================================

# 1 TITLE
s = new_slide()
kicker(s, "INDABA 2026  ·  3-HOUR CODE-ALONG", fill=GOLD, x=0.7, y=0.7)
text(s, 0.7, 1.55, 11.9, 2.2,
     [[("Automating Data Pipelines", {})],
      [("with ", {"color": TEXT}), ("Python", {"color": GREEN})]],
     size=54, bold=True, color=TEXT, font=HEAD, spacing=1.0)
text(s, 0.72, 3.62, 10.9, 0.8,
     "Take a year of messy agrivoltaics sensor logs and turn them into one "
     "trustworthy answer — with a single command.", size=17, color=MUTED,
     font=BODY, spacing=1.2)
step_flow(s, 0.7, 4.85, 11.9,
          [("01", "LOAD"), ("02", "RESHAPE"), ("03", "CONTRACT"),
           ("04", "CLEAN"), ("05", "VALIDATE"), ("06", "INSIGHTS")],
          chip_h=0.66, color=GOLD)
text(s, 0.7, 6.7, 11.9, 0.4,
     "Agrivoltaic Dataset Ghana  ·  Responsible AI Lab, KNUST  ·  FAIR Forward  ·  CC-BY 4.0",
     size=10.5, color=MUTED, font=BODY)

# 2 MEET AMA
s = new_slide()
kicker(s, "THE STORY", fill=GOLD)
title_block(s, "Meet Ama and her agrivoltaics question",
            "A real research question, buried inside a pile of messy spreadsheets.")
rrect(s, 0.7, 2.5, 4.6, 4.0, fill=PANEL, radius=0.07)
oval(s, 1.0, 2.85, 0.85, GREEN_D)
text(s, 1.0, 2.85, 0.85, 0.85, "A", size=30, bold=True, color=TEXT, font=HEAD,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 2.05, 2.92, 3.1, 0.5, "Ama", size=24, bold=True, color=TEXT, font=HEAD)
text(s, 2.07, 3.42, 3.1, 0.4, "KNUST researcher · Kumasi, Ghana", size=12.5,
     color=MUTED, font=BODY)
text(s, 1.0, 4.05, 4.0, 0.4,
     [[("HER QUESTION", {"color": GOLD, "bold": True, "size": 11, "spacing": 1.2})]],
     font=HEAD)
text(s, 1.0, 4.42, 4.0, 1.8,
     "Do the raised solar panels keep the crops cooler — while still "
     "generating power?", size=17, color=TEXT, font=BODY, spacing=1.2, italic=True)
text(s, 5.65, 2.55, 7, 0.4,
     [[("WHAT SHE ACTUALLY HAS", {"color": MUTED, "bold": True, "size": 11.5,
        "spacing": 1.2})]], font=HEAD)
step_flow(s, 5.65, 2.95, 6.95,
          [("01", "5 FILES"), ("02", "~150 SHEETS"), ("03", "5-MIN LOGS")],
          chip_h=0.6, color=GREEN)
rrect(s, 5.65, 4.05, 6.95, 2.45, fill=PANEL, radius=0.07)
oval(s, 5.95, 4.38, 0.16, CORAL)
text(s, 6.25, 4.3, 6, 0.4, "The risk", size=16, bold=True, color=TEXT, font=HEAD)
text(s, 5.95, 4.95, 6.4, 1.4,
     "Five monthly Excel files, a sheet per day, sensor codes in the headers, "
     "and gaps everywhere. Read one wrong and the answer quietly changes.",
     size=15, color=MUTED, font=BODY, spacing=1.25)
footer(s, F + "  ·  OPENING", 2)

# 3 AUTOMATION != TRUST
s = new_slide()
kicker(s, "THE TENSION", fill=GOLD)
title_block(s, "Automation is not the same as trust",
            "A fast wrong answer is still wrong. We are building for the third column.")
concept_cards(s, [
    ("Manual", ["Steps live in someone's memory", "Fixes are invisible",
                "Failures surface late", "Handover is painful"], CORAL),
    ("Automated", ["Runs without re-clicking", "But assumptions stay hidden",
                   "Speed is not proof", "Wrong, just faster"], GOLD),
    ("Trustworthy", ["Stages are explicit", "Inputs are checked",
                     "Outputs are reproducible", "Failures are loud"], GREEN),
], y=2.55, h=3.4)
bottom_note(s, [("This workshop moves you ", {"color": MUTED}),
                ("left → right", {"color": GREEN, "bold": True}),
                (".", {"color": MUTED})])
footer(s, F + "  ·  OPENING", 3)

# 4 WHAT IS A PIPELINE
s = new_slide()
kicker(s, "MENTAL MODEL", fill=GREEN)
title_block(s, "So what is a data pipeline?",
            "A repeatable system that moves data through explicit, named stages.")
step_flow(s, 0.7, 2.75, 11.9,
          [("01", "SOURCE"), ("02", "INGEST"), ("03", "TRANSFORM"),
           ("04", "VALIDATE"), ("05", "PUBLISH")], chip_h=0.85, color=GREEN,
          label_size=13)
concept_cards(s, [
    ("Data", ["The values flowing", "through the system."], GREEN),
    ("Logic", ["The rules that change", "or reject those values."], GOLD),
    ("Control", ["The order, failure behavior,", "and evidence of each run."], CORAL),
], y=4.15, h=2.0, head_size=17)
footer(s, F + "  ·  OPENING", 4)

# 5 BIG IDEA: ONE COMMAND (real output)
s = new_slide()
kicker(s, "THE BIG IDEA", fill=GREEN)
title_block(s, "One command should rebuild everything",
            "If trust lives in your head, it dies when you close the laptop. Put it in code.")
code_panel(s, 0.7, 2.5, 6.9, 3.95, "bash",
           "$ python run_pipeline.py\n\n"
           "Starting Ghana agrivoltaics pipeline\n"
           "Reshaped 1,120,152 readings from raw\n"
           "Validation: 4 warnings recorded\n"
           "writing daily summary...        ok\n"
           "writing midday comparison...    ok\n"
           "writing microclimate chart...   ok\n\n"
           "Pipeline completed successfully", size=12.5)
card(s, 7.95, 2.5, 4.65, 1.83, "The command is boring",
     ["One line turns 5 messy workbooks", "into trusted tables and a chart."],
     accent=GOLD)
card(s, 7.95, 4.52, 4.65, 1.83, "The trust is the workshop",
     ["Every line above is a stage we", "build together so it holds up."],
     accent=GREEN)
footer(s, F + "  ·  MENTAL MODEL", 5)

# 6 THE 3-HOUR ROUTE
s = new_slide()
kicker(s, "THE MAP", fill=GOLD)
title_block(s, "Your next three hours",
            "Short concept bursts, live builds, hands-on exercises, and checkpoints.")
route = [("0:00", "Mental model", "25 min"), ("0:25", "Load raw sheets", "20 min"),
         ("0:45", "Reshape & profile", "25 min"), ("1:10", "Contract", "20 min"),
         ("1:30", "Break", "10 min"), ("1:40", "Clean", "25 min"),
         ("2:05", "Validate", "25 min"), ("2:30", "Insights", "25 min"),
         ("2:55", "Wrap-up", "5 min")]
col_w, rh, rgap = 5.75, 0.56, 0.12
for i, (t, lab, mins) in enumerate(route):
    col = i // 5
    row = i % 5
    rx = 0.7 + col * (col_w + 0.4)
    ry = 2.65 + row * (rh + rgap)
    accent = GOLD if lab == "Break" else GREEN
    rrect(s, rx, ry, col_w, rh, fill=PANEL, radius=0.12)
    text(s, rx + 0.25, ry, 0.95, rh, t, size=14, bold=True, color=accent,
         font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx + 1.35, ry, col_w - 2.4, rh, lab, size=15, bold=True, color=TEXT,
         font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    text(s, rx + col_w - 1.4, ry, 1.2, rh, mins, size=12, color=MUTED, font=BODY,
         anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
footer(s, F + "  ·  OPENING", 6)

# 7 PREDICT
s = new_slide()
kicker(s, "PREDICT  ·  COMMIT BEFORE WE INSPECT", fill=GOLD)
title_block(s, "Which change should break the pipeline?",
            "No peeking at the data yet — pick the one that must stop the run.")
opts = [("A", "A new monthly Excel arrives with the usual day-sheets.", LINE),
        ("B", "A sheet uses an unknown plot code like ZZ-XX.", GREEN),
        ("C", "One temperature cell reads 350 °C.", LINE),
        ("D", "A sheet puts its header on row 2 instead of row 1.", LINE)]
for i, (letter, body, accent) in enumerate(opts):
    col = i % 2; row = i // 2
    cx = 0.7 + col * 6.15
    cy = 2.6 + row * 1.45
    rrect(s, cx, cy, 5.75, 1.25, fill=PANEL, radius=0.08,
          line=(accent if accent != LINE else None), line_w=1.5)
    oval(s, cx + 0.3, cy + 0.34, 0.56, PANEL2)
    text(s, cx + 0.3, cy + 0.34, 0.56, 0.56, letter, size=22, bold=True,
         color=GOLD, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx + 1.1, cy, 4.4, 1.25, body, size=14.5, color=TEXT, font=BODY,
         anchor=MSO_ANCHOR.MIDDLE, spacing=1.1)
bottom_note(s, [("DISCUSS   ", {"color": GOLD, "bold": True, "size": 12}),
                ("An unknown plot code is a contract break — it must stop the run. "
                 "A 350° reading we quarantine; a shifted header we detect.",
                 {"color": TEXT})])
footer(s, F + "  ·  OPENING", 7)

# ============================================================================
# MODULE 1 — LOAD
# ============================================================================
M1 = F + "  ·  MODULE 1 · LOAD"
divider_slide("01", "MODULE 1  ·  20 MINUTES", "Load the Raw Sheets",
              "Before we reshape anything, prove the pipeline can find and open "
              "every day of data.",
              "every day-sheet in every monthly workbook, readable exactly as "
              "recorded.", F + "  ·  MODULE 1", 8)

# 9 RAW IS EVIDENCE
s = new_slide()
kicker(s, "PRINCIPLE", fill=GREEN)
title_block(s, "Raw data is evidence — don't overwrite it",
            "Keep the original Excel files exactly as received. Everything else is rebuildable.")
concept_cards(s, [
    ("data/raw/", ["The provider's Excel files", "Immutable in our workflow",
                   "Messy headers and all"], GREEN),
    ("data/processed/", ["Created by our code", "Can always be rebuilt",
                         "Carries documented assumptions"], GOLD),
    ("Why keep both?", ["Audit what changed", "Fix the reshape logic",
                        "Reproduce any past output"], CORAL),
], y=2.55, h=3.4)
bottom_note(s, [("Rule of thumb:   ", {"color": MUTED, "bold": True}),
                ("if a step can't be re-run from raw, it isn't part of the pipeline yet.",
                 {"color": TEXT})])
footer(s, M1, 9)

# 10 DATA SHAPE
s = new_slide()
kicker(s, "KNOW YOUR INPUTS", fill=GREEN)
title_block(s, "One workbook per month, one sheet per day",
            "The data's structure is the first thing to understand — before any code.")
concept_cards(s, [
    ("5 workbooks", ["One Excel file per month", "May, Jun, Jul, Aug, Oct 2024",
                     "(.xlsx, not CSV)"], GREEN),
    ("~30 sheets each", ["One sheet per day", "Named like '1 08 24'",
                         "≈ 150 day-sheets in all"], GOLD),
    ("5-minute readings", ["A row every 5 minutes", "06:30 → dusk, ~159 rows",
                           "1.1M readings total"], CORAL),
], y=2.55, h=3.4)
bottom_note(s, [("This isn't a tidy weekly CSV — ", {"color": MUTED}),
                ("it's high-frequency sensor telemetry, nested in spreadsheets.",
                 {"color": TEXT})])
footer(s, M1, 10)

# 11 LIVE CODE: find files + sheets
code_with_notes(11, "LIVE CODE  ·  src/ingest.py", GOLD,
    "Find the files, then list their days",
    "pathlib finds the workbooks; pandas lists the day-sheets inside each one.",
    "src/ingest.py",
    "SUPPORTED_SUFFIXES = {\".csv\", \".xls\", \".xlsx\"}\n\n"
    "def list_tabular_files(raw_dir):\n"
    "    raw_path = Path(raw_dir)\n"
    "    return sorted(\n"
    "        path\n"
    "        for path in raw_path.glob(\"*\")\n"
    "        if path.is_file()\n"
    "        and path.suffix.lower() in SUPPORTED_SUFFIXES\n"
    "    )\n\n"
    "def list_day_sheets(path):\n"
    "    return pd.ExcelFile(path).sheet_names",
    "WHY THIS WAY",
    [("Files are values", "Discovered, sorted, and testable — no hard-coded names."),
     ("Sheets are data too", "Each workbook hides ~30 days we must enumerate."),
     ("One allow-list", "Only .csv/.xls/.xlsx are ever considered.")],
    M1, code_w=7.35, csize=12.5)

# 12 LIVE CODE: read raw sheet
code_with_notes(12, "LIVE CODE  ·  src/ingest.py", GOLD,
    "Read each sheet exactly as recorded",
    "header=None means pandas assumes nothing — we keep every messy row.",
    "src/ingest.py",
    "def load_raw_sheet(path, sheet_name):\n"
    "    return pd.read_excel(\n"
    "        path, sheet_name=sheet_name, header=None\n"
    "    )\n\n"
    "def iter_raw_sheets(raw_dir):\n"
    "    for file_path in list_tabular_files(raw_dir):\n"
    "        for sheet in list_day_sheets(file_path):\n"
    "            yield file_path, sheet, load_raw_sheet(\n"
    "                file_path, sheet)",
    "DESIGN CHOICE",
    [("Assume nothing", "header=None keeps the two header rows as data for now."),
     ("Stream, don't hoard", "yield hands back one day at a time."),
     ("Loud on bad files", "An unreadable workbook raises, never guesses.")],
    M1, code_w=7.35, csize=12.5)

# 13 EXERCISE
exercise_slide("Load the raw sheets",
    "Work from evidence — open the files and see what's really there.", 7,
    [[("List every workbook in ", {}),
      ("data/raw/", {"font": MONO, "color": GREEN, "size": 14}), (".", {})],
     [("Print the ", {}),
      ("sheet names", {"font": MONO, "color": GREEN, "size": 14}),
      (" of the first file.", {})],
     [("Read one day-sheet with ", {}),
      ("header=None", {"font": MONO, "color": GREEN, "size": 13.5}),
      (" and look.", {})],
     [("Stretch: count total sheets across all files.", {})]],
    "src/ingest.py", M1, 13)

# 14 CHECKPOINT
checkpoint_slide("Inputs are under control",
    "Pause and check the result before we add another layer.",
    ["All five workbooks are listed.",
     "You can enumerate the day-sheets inside each file.",
     "A single raw sheet can be read and inspected.",
     "No source file was edited by hand."], M1, 14)

# ============================================================================
# MODULE 2 — RESHAPE & PROFILE
# ============================================================================
M2 = F + "  ·  MODULE 2 · RESHAPE & PROFILE"
divider_slide("02", "MODULE 2  ·  25 MINUTES", "Reshape & Profile",
              "Turn the wide, two-header sensor grid into tidy rows — then measure "
              "what you actually have.",
              "one tidy long table of readings, plus a profile of its coverage "
              "and quality.", F + "  ·  MODULE 2", 15)

# 16 THE TWO-HEADER PROBLEM
s = new_slide()
kicker(s, "THE REAL CHALLENGE", fill=GOLD)
title_block(s, "The raw sheet has two header rows",
            "Plot codes sit above measurement types, with time down the side.")
code_panel(s, 0.7, 2.5, 7.0, 3.9, "raw sheet · 1 08 24",
           "         AG-PV P3        AO-TI P1   ...\n"
           "         Irr     T       RH     T\n"
           "06:30    57      —       89.4   24.2\n"
           "06:35    69      —       87.8   24.8\n"
           "06:40    78      —       87.3   25.1\n"
           "  ...\n"
           "11:34  Need network restart - ETH...", size=12)
card(s, 7.95, 2.5, 4.65, 1.25, "Codes are merged",
     ["One plot code spans several", "columns; blanks must be filled."],
     accent=GOLD)
card(s, 7.95, 3.88, 4.65, 1.25, "Header row drifts",
     ["Some sheets add a blank top", "row — find it, don't assume it."],
     accent=CORAL)
card(s, 7.95, 5.26, 4.65, 1.14, "Cells lie",
     ["Junk text hides among times", "and numbers."], accent=GREEN)
footer(s, M2, 16)

# 17 TIDY DATA
s = new_slide()
kicker(s, "MENTAL MODEL", fill=GREEN)
title_block(s, "Tidy data: one row per observation",
            "Every later stage becomes easy once each reading is its own row.")
data_table(s, 0.7, 2.7, 11.9,
           ["date", "raw_time", "plot_code", "measurement", "value"],
           [["2024-08-01", "06:30:00", "AG-PV P3", "irradiance", "57"],
            ["2024-08-01", "06:30:00", "AO-TI P1", "temperature", "24.2"],
            ["2024-08-01", "06:35:00", "AG-PV P3", "irradiance", "69"]],
           col_fr=[2.4, 2.2, 2.4, 2.6, 1.6], row_h=0.56, fs=13)
bottom_note(s, [("Wide grid in, ", {"color": MUTED}),
                ("tidy long table out", {"color": GREEN, "bold": True}),
                (" — now grouping, filtering, and joining all just work.",
                 {"color": MUTED})], y=6.2)
footer(s, M2, 17)

# 18 LIVE CODE: detect header
code_with_notes(18, "LIVE CODE  ·  src/reshape.py", GOLD,
    "Find the header row instead of assuming it",
    "We locate the measurement row by its labels, so a drifting header can't fool us.",
    "src/reshape.py",
    "MEASUREMENT_LABELS = {\n"
    "    \"Irr (W/m2)\", \"T (oC)\",\n"
    "    \"RH (%)\", \"P (mm)\"}\n\n"
    "def detect_header_rows(raw):\n"
    "    for i in range(min(6, len(raw))):\n"
    "        labels = {str(v).strip()\n"
    "                  for v in raw.iloc[i]\n"
    "                  if isinstance(v, str)}\n"
    "        if labels & MEASUREMENT_LABELS:\n"
    "            return i - 1, i\n"
    "    raise ValueError(\"No header row found\")",
    "READ IT AS A SEARCH",
    [("Scan the top rows", "The measurement labels are the fingerprint we match."),
     ("Plot row is above", "Return both: plot-code row and measurement row."),
     ("Fail loudly", "A sheet with no header raises, never guesses.")],
    M2, code_w=7.55, csize=12.5)

# 19 LIVE CODE: tidy_sheet
code_with_notes(19, "LIVE CODE  ·  src/reshape.py", GOLD,
    "Reshape one day into tidy rows",
    "Forward-fill the merged plot codes, then emit one column of readings at a time.",
    "src/reshape.py",
    "def tidy_sheet(raw, sheet_date):\n"
    "    plot_row, meas_row = detect_header_rows(raw)\n"
    "    plots = raw.iloc[plot_row].ffill()\n"
    "    meas = raw.iloc[meas_row]\n"
    "    body = raw.iloc[meas_row + 1:]\n"
    "    time = body.iloc[:, 0]\n\n"
    "    columns = []\n"
    "    for c in range(1, raw.shape[1]):\n"
    "        p, m = plots.iloc[c], meas.iloc[c]\n"
    "        if m not in MEASUREMENT_LABELS:\n"
    "            continue\n"
    "        columns.append(pd.DataFrame({\n"
    "            \"date\": sheet_date, \"raw_time\": time,\n"
    "            \"plot_code\": p, \"measurement\": m,\n"
    "            \"value\": body.iloc[:, c]}))\n"
    "    return pd.concat(columns)",
    "THE KEY MOVE",
    [(".ffill() the codes", "Carries each plot code across its merged blanks."),
     ("Skip non-sensors", "Empty trailing columns never make it in."),
     ("Concat to long", "One tidy frame per column, stacked together.")],
    M2, code_w=7.75, csize=11.5)

# 20 LIVE CODE: profile
code_with_notes(20, "LIVE CODE  ·  src/profile_data.py", GOLD,
    "Profile the tidy readings",
    "One row per column turns a million readings into a readable diagnostic.",
    "src/profile_data.py",
    "def profile_table(df):\n"
    "    return pd.DataFrame({\n"
    "        \"column\": df.columns,\n"
    "        \"dtype\": [str(df[c].dtype)\n"
    "                  for c in df.columns],\n"
    "        \"missing_pct\": [\n"
    "            round(df[c].isna().mean() * 100, 2)\n"
    "            for c in df.columns],\n"
    "        \"unique_values\": [\n"
    "            df[c].nunique() for c in df.columns],\n"
    "    })",
    "READ IT AS QUESTIONS",
    [("What is each column?", "Name and dtype, straight from the data."),
     ("How much is missing?", "Sensor gaps show up as a missing percentage."),
     ("How many distinct?", "Cardinality hints at codes vs measurements.")],
    M2, code_w=7.35, csize=12.5)

# 21 READ THE PROFILE (real numbers)
s = new_slide()
kicker(s, "INSPECT", fill=GREEN)
title_block(s, "Read the real profile like a detective",
            "These are the actual numbers from 1.1M readings — each one is a question.")
data_table(s, 0.7, 2.62, 11.9,
           ["column", "dtype", "missing_pct", "unique_values"],
           [["value", "float64", "46.32", "3723"],
            ["raw_time", "object", "2.62", "424"],
            ["timestamp", "datetime64", "3.10", "31207"],
            ["plot_code", "str", "0.00", "17"],
            ["measurement", "str", "0.00", "4"]],
           col_fr=[3, 2.4, 2.2, 2.4], row_h=0.5, fs=13)
bottom_note(s, [("46% of values are missing — ", {"color": TEXT}),
                ("sensor gaps, not errors. That's a lead to understand, not erase.",
                 {"color": GOLD, "bold": True})], y=6.25)
footer(s, M2, 21)

# 22 EXERCISE
exercise_slide("Reshape and profile",
    "Get every day into one tidy table, then measure what you have.", 8,
    [[("Reshape one day-sheet with ", {}),
      ("tidy_sheet()", {"font": MONO, "color": GREEN, "size": 13.5}), (".", {})],
     [("Build the full long table across all files.", {})],
     [("Profile it — find the ", {}),
      ("missing_pct", {"font": MONO, "color": GREEN, "size": 13.5}),
      (" of value.", {})],
     [("Find one column whose cardinality surprises you.", {})]],
    "src/reshape.py · src/profile_data.py", M2, 22)

# 23 CHECKPOINT
checkpoint_slide("We can describe the data",
    "Pause and check the result before we add another layer.",
    ["Every day-sheet reshapes into tidy long rows.",
     "The full table holds all 1.1M readings.",
     "Missingness and cardinality are measured, not guessed.",
     "The raw evidence is still untouched."], M2, 23)

# ============================================================================
# MODULE 3 — CONTRACT
# ============================================================================
M3 = F + "  ·  MODULE 3 · CONTRACT"
divider_slide("03", "MODULE 3  ·  20 MINUTES", "Define the Data Contract",
              "The plot codes carry the whole experiment. Make that meaning explicit.",
              "a decoded vocabulary that turns cryptic codes into treatments, "
              "stations, and units.", F + "  ·  MODULE 3", 24)

# 25 THE CODES ARE THE EXPERIMENT
s = new_slide()
kicker(s, "MENTAL MODEL", fill=GREEN)
title_block(s, "The plot codes ARE the experiment",
            "Three plots answer Ama's question — and every code names one of them.")
concept_cards(s, [
    ("AO — Plot 1", ["Open control field", "No panels, full sun",
                     "The crop baseline"], GOLD),
    ("AG — Plot 2", ["Agrivoltaic system", "Raised panels + crops",
                     "The treatment we test"], GREEN),
    ("PO — Plot 3", ["Ground-mounted PV", "Bare land, energy only",
                     "Full-sun light reference"], CORAL),
], y=2.55, h=3.4)
bottom_note(s, [("Suffixes name the sensor:  ", {"color": MUTED}),
                ("PV", {"color": TEXT, "bold": True}),
                (" = irradiance · ", {"color": MUTED}),
                ("TI / SS", {"color": TEXT, "bold": True}),
                (" = temp & humidity stations · ", {"color": MUTED}),
                ("WS", {"color": TEXT, "bold": True}),
                (" = weather station", {"color": MUTED})])
footer(s, M3, 25)

# 26 SCHEMA VS CONTRACT
s = new_slide()
kicker(s, "CONCEPT", fill=GREEN)
title_block(s, "Schema is necessary, but not sufficient",
            "Structure says how data is represented; a contract says what it means.")
card(s, 0.7, 2.6, 5.75, 3.2, "Schema check  —  can we read it?",
     ["The value column is numeric.", "The plot_code column exists.",
      "The time string parses."], accent=GOLD, head_size=16)
card(s, 6.85, 2.6, 5.75, 3.2, "Contract check  —  can we trust it?",
     ["AG means agrivoltaic, always.", "Irradiance is in W/m², 0–1500.",
      "Each row is one timed reading."], accent=GREEN, head_size=16)
bottom_note(s, [("Schema answers ", {"color": MUTED}),
                ("“can we read it?”", {"color": GOLD, "bold": True}),
                ("   Contract answers ", {"color": MUTED}),
                ("“do we know what it means?”", {"color": GREEN, "bold": True})])
footer(s, M3, 26)

# 27 DECODE A PLOT CODE
s = new_slide()
kicker(s, "DESIGN", fill=GOLD)
title_block(s, "Decode one plot code into meaning",
            "Every code splits into a treatment, a sensor station, and a replicate.")
rrect(s, 0.7, 2.7, 4.2, 1.0, fill=PANEL2, radius=0.12, line=GREEN, line_w=1.0)
text(s, 0.7, 2.7, 4.2, 1.0, "AG-PV P3", size=26, bold=True, color=GREEN, font=MONO,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, 5.1, 3.0, color=GOLD, size=26)
rows = [("AG", "treatment", "agrivoltaic"), ("PV", "station", "irradiance sensor"),
        ("P3", "replicate", "plot 3 of 3")]
for i, (code, role, mean) in enumerate(rows):
    yy = 2.62 + i * 0.78
    rrect(s, 5.9, yy, 6.7, 0.66, fill=PANEL, radius=0.12)
    text(s, 6.1, yy, 1.2, 0.66, code, size=16, bold=True, color=GOLD, font=MONO,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 7.3, yy, 2.2, 0.66, role, size=13, color=MUTED, font=BODY,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.5, yy, 3.0, 0.66, mean, size=14.5, bold=True, color=TEXT, font=HEAD,
         anchor=MSO_ANCHOR.MIDDLE)
bottom_note(s, [("Decode once, in one place, ", {"color": MUTED}),
                ("and the whole pipeline speaks treatments instead of codes.",
                 {"color": TEXT})])
footer(s, M3, 27)

# 28 LIVE CODE: parse_plot_code
code_with_notes(28, "LIVE CODE  ·  src/contract.py", GOLD,
    "Encode the legend as code",
    "A small map and a split turn every code into structured, trustworthy fields.",
    "src/contract.py",
    "TREATMENT_BY_PREFIX = {\n"
    "    \"AG\": \"agrivoltaic\",\n"
    "    \"AO\": \"open_sun_control\",\n"
    "    \"PO\": \"ground_mounted_pv\",\n"
    "    \"WS\": \"ambient\",\n"
    "}\n\n"
    "def parse_plot_code(code):\n"
    "    base, _, replicate = code.partition(\" \")\n"
    "    prefix, _, station = base.partition(\"-\")\n"
    "    return {\n"
    "        \"treatment\": TREATMENT_BY_PREFIX.get(\n"
    "            prefix, \"unknown\"),\n"
    "        \"station\": station or None,\n"
    "        \"replicate\": replicate or None,\n"
    "    }",
    "ONE SOURCE OF TRUTH",
    [("The map is the contract", "Allowed prefixes live in exactly one place."),
     ("Unknown stays visible", "An unmapped prefix becomes 'unknown', not a crash."),
     ("Pure and testable", "Give it a string, get back meaning.")],
    M3, code_w=7.55, csize=12)

# 29 EXERCISE
exercise_slide("Draft the contract",
    "Turn the codes you profiled into a decoded, stable vocabulary.", 8,
    [[("Map each prefix (AG/AO/PO/WS) to a treatment.", {})],
     [("Run ", {}),
      ("parse_plot_code()", {"font": MONO, "color": GREEN, "size": 13.5}),
      (" on a few real codes.", {})],
     [("Apply the contract to add treatment + station.", {})],
     [("List the allowed range for each measurement.", {})]],
    "src/contract.py", M3, 29)

# 30 CHECKPOINT
checkpoint_slide("The pipeline has a vocabulary",
    "Pause and check the result before we add another layer.",
    ["Every plot prefix maps to a known treatment.",
     "Codes decode into treatment, station, and replicate.",
     "Each measurement has a documented unit and range.",
     "Unknown codes are surfaced, not silently dropped."], M3, 30)

# ============================================================================
# BREAK
# ============================================================================
s = new_slide(PANEL)
text(s, 0.7, 1.55, 12, 1.6, "10", size=150, bold=True, color=GOLD, font=HEAD,
     align=PP_ALIGN.CENTER)
text(s, 0.7, 3.95, 12, 0.6, "MINUTE BREAK", size=30, bold=True, color=TEXT,
     font=HEAD, align=PP_ALIGN.CENTER)
text(s, 1.7, 4.75, 10, 0.8,
     "When we're back: clean the values, validate the ranges, and break the "
     "pipeline on purpose.", size=17, color=MUTED, font=BODY, align=PP_ALIGN.CENTER,
     spacing=1.25)
text(s, 0.7, 5.85, 12, 0.5,
     [[("CATCH-UP CHECK     ", {"color": GREEN, "bold": True, "size": 12, "spacing": 1}),
       ("files loaded   ·   reshaped to tidy   ·   contract decoded",
        {"color": MUTED, "size": 13})]], align=PP_ALIGN.CENTER)
footer(s, F + "  ·  HALFWAY", 31)

# ============================================================================
# MODULE 4 — CLEAN
# ============================================================================
M4 = F + "  ·  MODULE 4 · CLEAN"
divider_slide("04", "MODULE 4  ·  25 MINUTES", "Clean and Standardize",
              "Turn messy cells into trustworthy timestamps and numbers — without "
              "erasing evidence.",
              "real timestamps, numeric values, and standardized treatment labels.",
              F + "  ·  MODULE 4", 32)

# 33 DETERMINISTIC & IDEMPOTENT
s = new_slide()
kicker(s, "TWO PROPERTIES", fill=GREEN)
title_block(s, "Make every transform safe to re-run",
            "Two properties keep a re-run from quietly changing your results.")
card(s, 0.7, 2.6, 5.75, 2.5, "Deterministic",
     ["Same input + same code", "→ always the same output.",
      "No hidden manual step."], accent=GREEN, head_size=17)
card(s, 6.85, 2.6, 5.75, 2.5, "Idempotent",
     ["Apply it again and again", "→ the state stops changing.",
      "Safe retries, no drift."], accent=GOLD, head_size=17)
rrect(s, 0.7, 5.3, 11.9, 0.95, fill=CODE_BG, radius=0.07, line=LINE)
text(s, 1.0, 5.3, 11.3, 0.95,
     [[("to_datetime(to_datetime(x))  ", {"font": MONO, "color": C_DEF, "size": 14}),
       ("→ the same timestamp, every time", {"color": MUTED, "size": 13.5})]],
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, M4, 33)

# 34 LIVE CODE: timestamps
code_with_notes(34, "LIVE CODE  ·  src/transform.py", GOLD,
    "Build real timestamps from messy cells",
    "Combine the sheet's date with each row's time; junk simply becomes NaT.",
    "src/transform.py",
    "cleaned[\"date\"] = pd.to_datetime(\n"
    "    cleaned[\"date\"], errors=\"coerce\")\n\n"
    "stamp = (\n"
    "    cleaned[\"date\"].dt.strftime(\"%Y-%m-%d\")\n"
    "    + \" \"\n"
    "    + cleaned[\"raw_time\"].astype(str)\n"
    ")\n"
    "cleaned[\"timestamp\"] = pd.to_datetime(\n"
    "    stamp, errors=\"coerce\")\n\n"
    "# '11:34 - Need network restart' -> NaT",
    "COERCE → EXPOSE",
    [("Don't crash", "errors='coerce' keeps the pipeline running on junk."),
     ("Don't hide", "Unparseable times become a visible NaT."),
     ("Date from the sheet", "The day comes from the sheet name, the time from the row.")],
    M4, code_w=7.35, csize=12.5)

# 35 LIVE CODE: coerce numerics
code_with_notes(35, "LIVE CODE  ·  src/transform.py", GOLD,
    "Coerce values; quarantine the impossible",
    "Make every reading numeric, then blank out anything physically impossible.",
    "src/transform.py",
    "cleaned[\"value\"] = pd.to_numeric(\n"
    "    cleaned[\"value\"], errors=\"coerce\")\n\n"
    "def quarantine_out_of_range(cleaned):\n"
    "    out = cleaned.copy()\n"
    "    for name, (low, high) in RANGES.items():\n"
    "        mask = out[\"measurement\"] == name\n"
    "        bad = mask & (\n"
    "            (out[\"value\"] < low)\n"
    "            | (out[\"value\"] > high))\n"
    "        out.loc[bad, \"value\"] = float(\"nan\")\n"
    "    return out",
    "KEEP THE ROW",
    [("Numeric or NaN", "'not recorded' coerces to NaN, never a fake number."),
     ("Quarantine ≠ delete", "A 350° reading is blanked; its row survives."),
     ("Evidence preserved", "We can still see when and where it happened.")],
    M4, code_w=7.55, csize=12)

# 36 LIVE CODE: standardize labels
code_with_notes(36, "LIVE CODE  ·  src/contract.py", GOLD,
    "Standardize labels through the contract",
    "One pass attaches the treatment, station, and unit to every reading.",
    "src/contract.py",
    "def apply_contract(long_table):\n"
    "    out = long_table.copy()\n"
    "    parts = out[\"plot_code\"].map(parse_plot_code)\n"
    "    out[\"treatment\"] = [p[\"treatment\"]\n"
    "                        for p in parts]\n"
    "    out[\"station\"] = [p[\"station\"]\n"
    "                      for p in parts]\n"
    "    out[\"measurement\"] = out[\"measurement\"].map(\n"
    "        MEASUREMENT_NAMES)\n"
    "    out[\"unit\"] = out[\"measurement\"].map(\n"
    "        MEASUREMENT_UNITS)\n"
    "    return out",
    "ONE VOCABULARY",
    [("Codes → treatments", "Downstream code says 'agrivoltaic', not 'AG'."),
     ("Names → units", "'Irr (W/m2)' becomes 'irradiance' + 'W/m2'."),
     ("Copy, don't mutate", ".copy() keeps the reshaped table intact.")],
    M4, code_w=7.55, csize=12)

# 37 EXERCISE
exercise_slide("Clean the readings",
    "Make times real, values numeric, and labels consistent.", 8,
    [[("Build a ", {}),
      ("timestamp", {"font": MONO, "color": GREEN, "size": 14}),
      (" from date + raw_time.", {})],
     [("Coerce ", {}), ("value", {"font": MONO, "color": GREEN, "size": 14}),
      (" to numeric with coerce.", {})],
     [("Apply the contract to label every row.", {})],
     [("Count how many timestamps became NaT.", {})]],
    "src/transform.py", M4, 37)

# 38 CHECKPOINT
checkpoint_slide("Messy cells are now trustworthy",
    "Pause and check the result before we add another layer.",
    ["Timestamps are real datetimes or visible NaT.",
     "Values are numeric or visible NaN.",
     "Every row carries its treatment, station, and unit.",
     "The reshaped table was copied, not mutated."], M4, 38)

# ============================================================================
# MODULE 5 — VALIDATE
# ============================================================================
M5 = F + "  ·  MODULE 5 · VALIDATE"
divider_slide("05", "MODULE 5  ·  25 MINUTES", "Validate Before Trusting",
              "Turn the contract into executable checks that decide what may be "
              "published.",
              "a validation report that stops structural failures and flags "
              "data-quality issues.", F + "  ·  MODULE 5", 39)

# 40 CLEAN VS VALIDATE
s = new_slide()
kicker(s, "DIFFERENT JOBS", fill=GREEN)
title_block(s, "Cleaning changes; validation judges",
            "A rule normalises a known variation. A check exposes an unknown violation.")
card(s, 0.7, 2.6, 5.75, 3.2, "Clean  —  change representation",
     ["'11:34 junk'  →  NaT", "'not recorded'  →  NaN",
      "'AG'  →  'agrivoltaic'"], accent=GOLD, head_size=16, body_size=14)
card(s, 6.85, 2.6, 5.75, 3.2, "Validate  —  judge acceptability",
     ["plot prefix is known", "irradiance sits in 0–1500",
      "humidity sits in 0–100"], accent=GREEN, head_size=16, body_size=14)
bottom_note(s, [("If cleaning could fix it, clean. If trust depends on it, ",
                 {"color": MUTED}), ("validate.", {"color": GREEN, "bold": True})])
footer(s, M5, 40)

# 41 TWO TIERS OF FAILURE
s = new_slide()
kicker(s, "DESIGN THE FAILURE PATH", fill=GOLD)
title_block(s, "Not every problem should stop the run",
            "Separate structural failures from the noise real sensors always produce.")
card(s, 0.7, 2.6, 5.75, 3.2, "Structural  →  STOP",
     ["Unknown plot code", "Missing required column",
      "The data isn't what we agreed.", "Save the report and halt."],
     accent=CORAL, head_size=16, body_size=14)
card(s, 6.85, 2.6, 5.75, 3.2, "Data quality  →  QUARANTINE",
     ["A few out-of-range spikes", "Some unparseable timestamps",
      "Expected from real sensors.", "Blank them, warn, carry on."],
     accent=GREEN, head_size=16, body_size=14)
bottom_note(s, [("A pipeline that halts on every glitch never ships; ", {"color": MUTED}),
                ("one that ignores structure can't be trusted.", {"color": TEXT})])
footer(s, M5, 41)

# 42 LIVE CODE: validate_ranges
code_with_notes(42, "LIVE CODE  ·  src/validate_data.py", GOLD,
    "Write small, composable checks",
    "Each check answers one question and returns evidence — it never stops the run itself.",
    "src/validate_data.py",
    "def validate_ranges(cleaned):\n"
    "    errors = []\n"
    "    for name, (low, high) in RANGES.items():\n"
    "        values = cleaned.loc[\n"
    "            cleaned[\"measurement\"] == name,\n"
    "            \"value\"]\n"
    "        bad = values[(values < low)\n"
    "                     | (values > high)]\n"
    "        if not bad.empty:\n"
    "            errors.append(\n"
    "                f\"{name}: {len(bad)} values \"\n"
    "                f\"outside [{low}, {high}]\")\n"
    "    return errors",
    "COMPOSABLE BY DESIGN",
    [("No printing, no exit", "It returns messages; the pipeline decides."),
     ("Returns evidence", "Counts and bounds, ready for one combined report."),
     ("Many tiny checks", "Ranges, known plots, timestamps — all roll up.")],
    M5, code_w=7.55, csize=12)

# 43 BREAK IT ON PURPOSE (real report)
s = new_slide()
kicker(s, "BREAK IT ON PURPOSE", fill=GOLD)
title_block(s, "A trustworthy pipeline fails loudly",
            "This is the real report from our run — every issue named, with evidence.")
code_panel(s, 0.7, 2.5, 6.9, 3.75, "outputs/reports/validation_report.csv",
           "status,message\n"
           "warning,\"temperature: 17766 values\n"
           "         outside [5, 60]\"\n"
           "warning,\"humidity: 445 values\n"
           "         outside [0, 100]\"\n"
           "warning,\"34778 rows have an\n"
           "         unparseable timestamp\"", size=12)
card(s, 7.95, 2.5, 4.65, 1.78, "Warnings → quarantine",
     ["Real sensors spike and drop.", "Blank the bad, keep the good."],
     accent=GREEN)
card(s, 7.95, 4.47, 4.65, 1.78, "Add an unknown plot → STOP",
     ["Inject a ZZ-XX code and the", "run halts before publishing."],
     accent=CORAL)
footer(s, M5, 43)

# 44 EXERCISE
exercise_slide("Break it on purpose",
    "Prove that data you can't trust never reaches a published table.", 8,
    [[("Add an unknown plot code and run the check.", {})],
     [("Confirm the pipeline stops before publishing.", {})],
     [("Add a 350 °C reading; watch it get quarantined.", {})],
     [("Open the saved ", {}),
      ("validation_report.csv", {"font": MONO, "color": GREEN, "size": 12.5}),
      (".", {})]],
    "src/validate_data.py", M5, 44)

# 45 CHECKPOINT
checkpoint_slide("Untrusted data cannot publish",
    "Pause and check the result before we add another layer.",
    ["Unknown plot codes stop the run.",
     "Out-of-range readings are quarantined, not published.",
     "Unparseable timestamps are counted and reported.",
     "Every run leaves a saved validation report."], M5, 45)

# ============================================================================
# MODULE 6 — INSIGHTS
# ============================================================================
M6 = F + "  ·  MODULE 6 · INSIGHTS"
divider_slide("06", "MODULE 6  ·  25 MINUTES", "Transform Into Insights",
              "Design the summary backward from Ama's question, then answer it.",
              "daily summaries, a midday microclimate comparison, and the chart "
              "that answers the question.", F + "  ·  MODULE 6", 46)

# 47 GRAIN / MIDDAY
s = new_slide()
kicker(s, "MENTAL MODEL", fill=GREEN)
title_block(s, "Choose the grain that answers the question",
            "1.1M raw readings can't be read — so we summarise to a meaningful grain.")
card(s, 0.7, 2.6, 5.75, 2.9, "Raw reading grain",
     ["One sensor, one plot,", "one 5-minute moment.",
      "Far too fine to compare."], accent=GOLD, head_size=16, body_size=14)
card(s, 6.85, 2.6, 5.75, 2.9, "Comparison grain",
     ["One treatment × measurement,", "averaged over the midday",
      "window (11:00–14:00)."], accent=GREEN, head_size=16, body_size=14)
bottom_note(s, [("We compare at midday — ", {"color": MUTED}),
                ("when sun, heat, and the panels' effect are all at their peak.",
                 {"color": TEXT})])
footer(s, M6, 47)

# 48 LIVE CODE: daily summary
code_with_notes(48, "LIVE CODE  ·  src/transform.py", GOLD,
    "Aggregate at the grain you declared",
    "Group by the dimensions; compute measures with named aggregations.",
    "src/transform.py",
    "return (\n"
    "    cleaned.dropna(\n"
    "        subset=[\"timestamp\", \"value\"])\n"
    "    .groupby([\"date\", \"treatment\",\n"
    "              \"station\", \"measurement\"],\n"
    "             dropna=False)\n"
    "    .agg(\n"
    "        observations=(\"value\", \"count\"),\n"
    "        mean_value=(\"value\", \"mean\"),\n"
    "        min_value=(\"value\", \"min\"),\n"
    "        max_value=(\"value\", \"max\"),\n"
    "    )\n"
    "    .reset_index()\n"
    ")",
    "READ THE QUERY",
    [("Dimensions", "date + treatment + station + measurement."),
     ("Measures", "Count and mean/min/max of the readings."),
     ("Drop the unusable", "Rows with no timestamp or value don't count.")],
    M6, code_w=7.55, csize=12)

# 49 LIVE CODE: midday comparison
code_with_notes(49, "LIVE CODE  ·  src/transform.py", GOLD,
    "Compare against the right baseline",
    "Each measurement is judged against the open reference that actually exists.",
    "src/transform.py",
    "BASELINE_BY_MEASUREMENT = {\n"
    "    \"temperature\": \"open_sun_control\",\n"
    "    \"humidity\": \"open_sun_control\",\n"
    "    \"irradiance\": \"ground_mounted_pv\",\n"
    "}\n\n"
    "for measurement, baseline in BASELINE...items():\n"
    "    ag = means.get((measurement, \"agrivoltaic\"))\n"
    "    ref = means.get((measurement, baseline))\n"
    "    diff = ag - ref\n"
    "    rows.append({\n"
    "        \"measurement\": measurement,\n"
    "        \"agrivoltaic\": round(ag, 2),\n"
    "        \"difference_pct\":\n"
    "            round(diff / ref * 100, 2)})",
    "THE BASELINE IS THE METRIC",
    [("Temp & humidity", "Compared to the open control field (AO)."),
     ("Irradiance", "No sensor in the field — use full-sun PV (PO)."),
     ("Honest by design", "We never compare against a baseline that isn't there.")],
    M6, code_w=7.75, csize=11.5)

# 50 THE RESULT (real numbers)
s = new_slide()
kicker(s, "THE ANSWER", fill=GREEN)
title_block(s, "Do the panels keep the crops cooler?",
            "Midday, agrivoltaic vs open control — straight from the pipeline output.")
stats = [("-7.3°C", "cooler at midday", "27.9° vs 35.1° · temperature", GREEN),
         ("-8.7 pts", "lower humidity", "15.6% vs 24.3% · drier air", GOLD),
         ("590", "W/m² on the panels", "energy generation continues", CORAL)]
cw, gap = (11.9 - 2 * 0.35) / 3, 0.35
for i, (big, lab, sub, accent) in enumerate(stats):
    cx = 0.7 + i * (cw + gap)
    rrect(s, cx, 2.65, cw, 2.9, fill=PANEL, radius=0.08)
    text(s, cx + 0.3, 3.0, cw - 0.6, 1.0, big, size=46, bold=True, color=accent,
         font=HEAD, align=PP_ALIGN.CENTER)
    text(s, cx + 0.3, 4.05, cw - 0.6, 0.4, lab, size=16, bold=True, color=TEXT,
         font=HEAD, align=PP_ALIGN.CENTER)
    text(s, cx + 0.3, 4.5, cw - 0.6, 0.8, sub, size=12.5, color=MUTED, font=BODY,
         align=PP_ALIGN.CENTER, spacing=1.15)
bottom_note(s, [("The dual-use win:  ", {"color": GREEN, "bold": True}),
                ("a noticeably cooler crop microclimate, with the panels still "
                 "capturing strong midday sun.", {"color": TEXT})])
footer(s, M6, 50)

# 51 EXERCISE
exercise_slide("Answer Ama's question",
    "Summarise to the comparison grain and produce the chart.", 8,
    [[("Build the daily plot summary.", {})],
     [("Average the midday window per treatment.", {})],
     [("Compute agrivoltaic vs its open reference.", {})],
     [("Save the comparison CSV and the chart.", {})]],
    "src/transform.py · run_pipeline.py", M6, 51)

# 52 CHECKPOINT
checkpoint_slide("You answered the question",
    "Pause and check the result before we wrap up.",
    ["Readings are summarised to a daily, per-plot grain.",
     "The midday comparison uses the correct baseline per measurement.",
     "The temperature difference is real and reproducible.",
     "Processed tables and a chart are saved to disk."], M6, 52)

# ============================================================================
# WRAP
# ============================================================================
W = F + "  ·  WRAP"

# 53 REPRODUCIBLE != CORRECT
s = new_slide()
kicker(s, "DISCUSSION", fill=GOLD)
title_block(s, "A reproducible result can still mislead",
            "Trust includes responsible interpretation — not just a clean exit code.")
card(s, 0.7, 2.6, 5.75, 3.2, "The pipeline can say",
     ["Midday is cooler under the",
      "panels across this dataset,",
      "under the cleaning and",
      "aggregation rules we encoded."], accent=GREEN,
     head_size=16, body_size=14)
card(s, 6.85, 2.6, 5.75, 3.2, "The pipeline cannot say",
     ["That cooler air means higher",
      "yield, or that this holds across",
      "every season, soil, crop, and",
      "site in Ghana."], accent=CORAL,
     head_size=16, body_size=14)
footer(s, W, 53)

# 54 CARRY IT FORWARD
s = new_slide(PANEL)
kicker(s, "CARRY IT FORWARD", fill=GREEN, x=0.7, y=0.7)
text(s, 0.7, 1.25, 11.9, 1.0, "Every result travels with its context",
     size=36, bold=True, color=TEXT, font=HEAD)
text(s, 0.72, 2.15, 11, 0.5,
     "Whenever you hand off a number, hand off these alongside it.", size=16,
     color=MUTED, font=BODY)
items = ["Source", "Licence", "Grain", "Baseline", "Validation status",
         "Limitations"]
cw, gap = (11.9 - 2 * 0.35) / 3, 0.35
for i, it in enumerate(items):
    col = i % 3; row = i // 3
    cx = 0.7 + col * (cw + gap)
    cy = 3.15 + row * 1.35
    rrect(s, cx, cy, cw, 1.1, fill=BG, radius=0.1, line=LINE)
    oval(s, cx + 0.3, cy + 0.45, 0.18, GREEN)
    text(s, cx + 0.65, cy, cw - 0.8, 1.1, it, size=18, bold=True, color=TEXT,
         font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
footer(s, W, 54)

# 55 FINAL
s = new_slide()
kicker(s, "WORKSHOP COMPLETE", fill=GOLD, x=0.7, y=0.8)
text(s, 0.7, 1.4, 11.9, 1.6,
     [[("You built more than a script", {})]], size=46, bold=True, color=TEXT,
     font=HEAD)
text(s, 0.72, 2.45, 11, 0.5, "You learned how trustworthy pipelines think.",
     size=18, color=GREEN, font=BODY)
principles = ["Explicit inputs", "Tidy reshaping", "Decoded contracts",
              "Honest cleaning", "Executable checks", "Reproducible answers"]
cw, gap = (11.9 - 2 * 0.3) / 3, 0.3
for i, p in enumerate(principles):
    col = i % 3; row = i // 3
    cx = 0.7 + col * (cw + gap)
    cy = 3.4 + row * 0.92
    rrect(s, cx, cy, cw, 0.72, fill=PANEL, radius=0.14)
    text(s, cx, cy, cw, 0.72, p, size=15, bold=True, color=TEXT, font=HEAD,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rrect(s, 0.7, 5.55, 11.9, 0.75, fill=CODE_BG, radius=0.1, line=GREEN)
text(s, 1.0, 5.55, 11.3, 0.75,
     [[("THE TEST    ", {"color": GREEN, "bold": True, "size": 12.5}),
       ("Can someone else rerun it, understand what happened, and trust the result?",
        {"color": TEXT, "size": 15})]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 6.55, 11.9, 0.4,
     "Agrivoltaic Dataset Ghana  ·  Responsible AI Lab, KNUST  ·  FAIR Forward  ·  CC-BY 4.0",
     size=10.5, color=MUTED, font=BODY)
footer(s, W, 55)

prs.save("Workshop_Slide_Deck.pptx")
print("Saved Workshop_Slide_Deck.pptx with", len(prs.slides._sldIdLst), "slides")
